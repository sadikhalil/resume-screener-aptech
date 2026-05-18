from pptx import Presentation

def parse_ppt(file_path: str) -> str:
    """Extract plain text from a PPT/PPTX file slide by slide."""
    prs = Presentation(file_path)
    parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)

    return "\n".join(parts).strip()
